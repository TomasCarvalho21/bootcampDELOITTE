from src.ingestion.loaders.loaderBase import LoaderBase
from pptx import Presentation
import io

class LoaderPPTX(LoaderBase):

    def __init__(self,filepath:str):
        self.filepath=filepath
    
    def extract_metadata(self):
        presentation = Presentation(self.filepath)
        properties = presentation.core_properties
        
        metadata = {
            "author": properties.author,
            "title": properties.title,
            "subject": properties.subject,
            "keywords": properties.keywords,
            "comments": properties.comments,
            "last_modified_by": properties.last_modified_by,
            "last_printed": properties.last_printed,
            "revision": properties.revision,
            "created": properties.created,
            "modified": properties.modified,
        }
        
        self.metadata=metadata
        
        return self.metadata if self.all_keys_have_values(metadata=self.metadata) else False

    def extract_text(self):
        full_pptx_text=""
        presentation = Presentation(self.filepath)
        slides_text=self.extract_slide_text_without_single_keywords(presentation,2)
        notes_text=self.extract_ppt_notes(presentation)

        for slide,note in zip(slides_text,notes_text):
            full_pptx_text=full_pptx_text+slide+"\n"+note+"\n\n"
            
        return full_pptx_text

    
    def extract_slide_text_without_single_keywords(self,presentation:object, isolated_words:int):
        slides_text = []

        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if len(text.split()) > isolated_words:
                        slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))
        
        return slides_text

    def extract_ppt_notes(self,presentation:object):
        notes_text = []

        for slide in presentation.slides:
            notes_slide = slide.notes_slide
            if notes_slide:
                notes_text.append(notes_slide.notes_text_frame.text)
        
        return notes_text

    def all_keys_have_values(self, metadata:dict, value_check=lambda x: x is not None and x != ''):
        return all(value_check(value) for value in metadata.values())